import streamlit as st
import os
import shutil
import re
import uuid
import requests
from PIL import Image
import fitz
from docx import Document
from pptx import Presentation
from qdrant_client import models, QdrantClient
from sentence_transformers import SentenceTransformer
from qdrant_client.models import PointStruct

# Access secrets via st.secrets
QDRANT_URL = st.secrets["QDRANT_API_URL"]
QDRANT_API_KEY = st.secrets["QDRANT_API_KEY"]
api_key = st.secrets["GEMINI_API_KEY"]
model_name = "gemini-1.5-flash"

def generate_response(api_key, model_name, context, query):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
    context_text = "\n".join(context)
    headers = {"Content-Type": "application/json"}
    prompt_text = f"""
    You are a helpful AI assistant. Use the following context to answer the question precisely and in detail in points.
    
    Context: {context_text}
    
    Question: {query}
    
    Answer:
    """
    payload = {"contents": [{"parts": [{"text": prompt_text}]}]}
    try:
        response = requests.post(url, headers=headers, json=payload)
        if response.ok:
            response_json = response.json()
            if 'candidates' in response_json and response_json['candidates']:
                return response_json['candidates'][0]['content']['parts'][0]['text']
            else:
                return "No candidates found in the response."
        else:
            return f"Error: {response.status_code}\n{response.text}"
    except Exception as e:
        return f"An error occurred: {str(e)}"

def process_documents(input_directory, image_output_dir="extracted_images"):
    if os.path.exists(image_output_dir):
        shutil.rmtree(image_output_dir)
    os.makedirs(image_output_dir, exist_ok=True)
    extracted_data = {}
    for root, _, files in os.walk(input_directory):
        for file in files:
            file_path = os.path.join(root, file)
            file_name, file_ext = os.path.splitext(file)
            if file_ext.lower() == '.pdf':
                text = extract_text_and_images_from_pdf(file_path, file_name, image_output_dir)
                extracted_data[file_name] = text
            elif file_ext.lower() == '.docx':
                text = extract_text_and_images_from_word(file_path, file_name, image_output_dir)
                extracted_data[file_name] = text
            elif file_ext.lower() == '.pptx':
                text = extract_text_and_images_from_ppt(file_path, file_name, image_output_dir)
                extracted_data[file_name] = text
    return extracted_data

def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\x00-\x7F]+', '', text)
    return text.strip()

def extract_text_and_images_from_pdf(pdf_path, file_name, image_output_dir):
    extracted_text = ""
    with fitz.open(pdf_path) as pdf:
        for page_number, page in enumerate(pdf):
            text = page.get_text()
            extracted_text += clean_text(text) + " "
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = os.path.join(image_output_dir, f"{file_name}_page{page_number+1}_img{img_index+1}.{image_ext}")
                with open(image_filename, "wb") as image_file:
                    image_file.write(image_bytes)
    return extracted_text

def extract_text_and_images_from_word(word_path, file_name, image_output_dir):
    extracted_text = ""
    doc = Document(word_path)
    for para in doc.paragraphs:
        extracted_text += clean_text(para.text) + " "
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_part = rel.target_part
            image_data = image_part.blob
            image_filename = os.path.join(image_output_dir, f"{file_name}_img{rel.rId}.png")
            with open(image_filename, "wb") as img_file:
                img_file.write(image_data)
    return extracted_text

def extract_text_and_images_from_ppt(ppt_path, file_name, image_output_dir):
    extracted_text = ""
    prs = Presentation(ppt_path)
    for slide_number, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += clean_text(shape.text) + " "
            if shape.shape_type == 13:
                image = shape.image
                image_bytes = image.blob
                image_filename = os.path.join(image_output_dir, f"{file_name}_slide{slide_number+1}_img{shape.shape_id}.png")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
    return extracted_text

def split_text(text, chunk_size=800, chunk_overlap=110):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - chunk_overlap
    return chunks

def main():
    UPLOAD_DIR = "input_files"
    image_output_directory = "extracted_images"
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    os.makedirs(image_output_directory, exist_ok=True)

    qdrant_client = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)
    encoder = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

    if 'session_id' not in st.session_state:
        st.session_state.session_id = str(uuid.uuid4())
    if 'messages' not in st.session_state:
        st.session_state.messages = []

    st.set_page_config(page_title="SamurAI - Document Q&A", page_icon="🥷", layout="wide", initial_sidebar_state="expanded")
    st.title("🥷SamurAI - Document Q&A Chat Bot")
    
    st.sidebar.header("Upload Documents")
    uploaded_files = st.sidebar.file_uploader("Choose documents", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

    if uploaded_files:
        for filename in os.listdir(UPLOAD_DIR):
            os.remove(os.path.join(UPLOAD_DIR, filename))
        for uploaded_file in uploaded_files:
            file_path = os.path.join(UPLOAD_DIR, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getvalue())
        st.sidebar.success(f"Uploaded {len(uploaded_files)} file(s)")
        with st.spinner("Processing documents..."):
            docs = process_documents(UPLOAD_DIR, image_output_directory)
            for file_name, text in docs.items():
                chunks = split_text(text)
                batch_points = []
                for i, chunk in enumerate(chunks):
                    embedding = encoder.encode(chunk)
                    batch_points.append(PointStruct(
                        id=len(batch_points), 
                        vector=embedding,
                        payload={
                            "text": chunk,
                            "document_name": file_name,
                            "session_id": st.session_state.session_id
                        }
                    ))
                for i in range(0, len(batch_points), 100):
                    batch = batch_points[i:i+100]
                    qdrant_client.upsert(collection_name="my_documents", points=batch)
        images = [os.path.join(image_output_directory, img) for img in os.listdir(image_output_directory)]
        if images:
            st.write("### Extracted Images")
            cols = st.columns(4)
            for i, img_path in enumerate(images):
                with cols[i % 4]:
                    st.image(img_path, use_container_width=True)

    st.write("### Ask Questions about Your Documents")
    user_question = st.text_input("Enter your question here")

    if st.button("Ask") and user_question:
        try:
            query_embedding = encoder.encode(user_question).tolist()
            session_id = st.session_state.session_id
            sessionFilter = models.Filter(must=[models.FieldCondition(key="session_id", match=models.MatchValue(value=session_id))])
            search_result = qdrant_client.search(
                collection_name="my_documents", query_vector=query_embedding, query_filter=sessionFilter, limit=8
            )
            retrieved_context = "\n".join([hit.payload['text'] for hit in search_result])
            bot_response = generate_response(api_key, model_name, retrieved_context, user_question)
            st.session_state.messages.append({'user': user_question, 'bot': bot_response})
        except Exception as e:
            st.error(f"Error processing query: {e}")

    st.write("### Chat History")
    for message in st.session_state.messages:
        st.markdown(f"""<div style=\"background-color: #282828; padding: 10px; margin-bottom: 10px; border-radius: 5px;\"><strong>You:</strong> {message['user']}</div>""", unsafe_allow_html=True)
        st.markdown(f"""<div style=\"background-color: #282828; padding: 10px; margin-bottom: 10px; border-radius: 5px;\"><strong>SamurAI:</strong> {message['bot']}</div>""", unsafe_allow_html=True)

if __name__ == "__main__":
    main()
